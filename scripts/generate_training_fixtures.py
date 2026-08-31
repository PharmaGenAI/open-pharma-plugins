#!/usr/bin/env python3
"""Generate sample training fixture documents for the field-training capability.

Creates two files in the output directory:
  - sample_product_message.pdf  (5 pages, 13 approved claims for fictional ONCORIX)
  - sample_training_deck.pptx   (6 slides with speaker notes)

Usage:
  python3 scripts/generate_training_fixtures.py
  python3 scripts/generate_training_fixtures.py --output-dir /tmp/fixtures
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parent.parent

DEFAULT_OUTPUT_DIR = (
    REPO_ROOT / "src" / "capabilities" / "field-training" / "open_pharma_plugins_field_training" / "fixtures"
)


# ---------------------------------------------------------------------------
# PDF generation (reportlab)
# ---------------------------------------------------------------------------


def _generate_pdf(output_path: Path) -> None:
    from reportlab.lib.colors import HexColor
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

    doc = SimpleDocTemplate(
        str(output_path),
        pagesize=letter,
        leftMargin=inch,
        rightMargin=inch,
        topMargin=inch,
        bottomMargin=inch,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "DocTitle",
        parent=styles["Title"],
        fontSize=22,
        spaceAfter=12,
        textColor=HexColor("#1a3c6e"),
    )
    subtitle_style = ParagraphStyle(
        "DocSubtitle",
        parent=styles["Normal"],
        fontSize=12,
        spaceAfter=6,
        textColor=HexColor("#555555"),
        alignment=1,
    )
    heading_style = ParagraphStyle(
        "SectionHeading",
        parent=styles["Heading2"],
        fontSize=16,
        spaceBefore=12,
        spaceAfter=8,
        textColor=HexColor("#1a3c6e"),
    )
    claim_style = ParagraphStyle(
        "Claim",
        parent=styles["Normal"],
        fontSize=10,
        spaceBefore=6,
        spaceAfter=6,
        leftIndent=18,
        bulletIndent=0,
        leading=14,
    )
    body_style = ParagraphStyle(
        "BodyText",
        parent=styles["Normal"],
        fontSize=10,
        spaceBefore=4,
        spaceAfter=4,
        leading=14,
    )
    footer_style = ParagraphStyle(
        "Footer",
        parent=styles["Normal"],
        fontSize=8,
        textColor=HexColor("#999999"),
        spaceBefore=24,
    )

    story: list = []

    # ── Page 1: Cover / Indication ──
    story.append(Spacer(1, 1.5 * inch))
    story.append(Paragraph("ONCORIX (rivolumab)", title_style))
    story.append(Paragraph("Approved Product Messages", title_style))
    story.append(Spacer(1, 0.3 * inch))
    story.append(Paragraph("For use by field medical and commercial teams only", subtitle_style))
    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph("Document ID: APM-2026-001 &nbsp;|&nbsp; Effective: January 2026", subtitle_style))
    story.append(Spacer(1, 0.5 * inch))
    story.append(Paragraph("Approved Indication", heading_style))
    story.append(
        Paragraph(
            "ONCORIX is indicated for the treatment of adult patients with unresectable or metastatic melanoma.",
            body_style,
        )
    )
    story.append(
        Paragraph(
            "CONFIDENTIAL — This document contains proprietary information. "
            "Do not distribute outside the organisation.",
            footer_style,
        )
    )
    story.append(PageBreak())

    # ── Page 2: Efficacy Claims ──
    story.append(Paragraph("Efficacy", heading_style))

    efficacy_claims = [
        (
            "Claim 1 (Overall Survival)",
            "In the MERIDIAN-301 trial, ONCORIX demonstrated a statistically significant "
            "improvement in overall survival (OS) vs standard of care (median OS 24.1 months "
            "vs 13.6 months; HR 0.68; 95% CI: 0.53–0.87; p=0.002).",
        ),
        (
            "Claim 2 (Response Rate)",
            "The objective response rate (ORR) was 42.3% (95% CI: 35.1–49.8) in the "
            "ONCORIX arm compared to 16.2% in the control arm.",
        ),
        (
            "Claim 3 (Progression-Free Survival)",
            "Median progression-free survival (PFS) was 11.2 months (95% CI: 8.9–14.1) "
            "with ONCORIX vs 5.4 months with standard of care.",
        ),
        (
            "Claim 4 (Durability)",
            "Durable responses were observed: 78% of responders maintained response at 12 months.",
        ),
        (
            "Claim 5 (Subgroups)",
            "Subgroup analyses showed consistent OS benefit across pre-specified subgroups "
            "including PD-L1 expression level, ECOG status, and prior therapy.",
        ),
    ]
    for label, text in efficacy_claims:
        story.append(Paragraph(f"<b>{label}:</b> {text}", claim_style))
    story.append(PageBreak())

    # ── Page 3: Safety Claims ──
    story.append(Paragraph("Safety", heading_style))

    safety_claims = [
        (
            "Claim 6 (Common AEs)",
            "The most common adverse reactions (≥20%) were fatigue (38%), rash (27%), "
            "diarrhoea (24%), and nausea (21%).",
        ),
        (
            "Claim 7 (Immune-Mediated AEs)",
            "Immune-mediated adverse events occurred in 31% of patients. Most were Grade 1-2 "
            "and manageable with established protocols.",
        ),
        (
            "Claim 8 (Grade 3-4 AEs)",
            "Grade 3-4 treatment-related adverse events occurred in 18% of patients in the "
            "ONCORIX arm vs 12% in the control arm.",
        ),
        (
            "Claim 9 (Discontinuation)",
            "Treatment discontinuation due to adverse events occurred in 9% of patients.",
        ),
        (
            "Claim 10 (Long-Term Safety)",
            "No new safety signals were identified in the 24-month follow-up analysis.",
        ),
    ]
    for label, text in safety_claims:
        story.append(Paragraph(f"<b>{label}:</b> {text}", claim_style))
    story.append(PageBreak())

    # ── Page 4: Dosing & Administration ──
    story.append(Paragraph("Dosing &amp; Administration", heading_style))

    dosing_claims = [
        (
            "Claim 11 (Recommended Dose)",
            "The recommended dose of ONCORIX is 200 mg administered as an intravenous "
            "infusion over 30 minutes every 3 weeks.",
        ),
        (
            "Claim 12 (Hepatic Impairment)",
            "No dose adjustment is required for patients with mild hepatic impairment "
            "(bilirubin ≤ ULN and AST > ULN, or bilirubin 1–1.5× ULN).",
        ),
        (
            "Claim 13 (Duration)",
            "Treatment should continue until disease progression, unacceptable toxicity, "
            "or up to 24 months in patients without disease progression.",
        ),
    ]
    for label, text in dosing_claims:
        story.append(Paragraph(f"<b>{label}:</b> {text}", claim_style))
    story.append(PageBreak())

    # ── Page 5: Approved Talking Points ──
    story.append(Paragraph("Approved Talking Points", heading_style))

    talking_points = [
        "When discussing ONCORIX with oncologists, focus on the significant OS benefit "
        "demonstrated in MERIDIAN-301 and the manageable safety profile.",
        "Always present efficacy and safety data together to ensure fair balance.",
        "Refer HCPs to the full prescribing information for complete safety data.",
    ]
    for i, tp in enumerate(talking_points, 1):
        story.append(Paragraph(f"<b>{i}.</b> {tp}", claim_style))

    story.append(Spacer(1, 0.5 * inch))
    story.append(
        Paragraph(
            "End of Approved Product Messages. APM-2026-001.",
            footer_style,
        )
    )

    doc.build(story)


# ---------------------------------------------------------------------------
# PPTX generation (python-pptx)
# ---------------------------------------------------------------------------


def _generate_pptx(output_path: Path) -> None:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def _add_slide(
        title: str,
        body_lines: list[str],
        notes: str,
        *,
        is_title_slide: bool = False,
        subtitle: str | None = None,
    ) -> None:
        if is_title_slide:
            layout = prs.slide_layouts[0]
        else:
            layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(32) if is_title_slide else Pt(28)

        if is_title_slide and subtitle:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = subtitle
                    for para in ph.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.CENTER
                        for run in para.runs:
                            run.font.size = Pt(18)
                    break

        if not is_title_slide and body_lines:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    tf = ph.text_frame
                    tf.clear()
                    for i, line in enumerate(body_lines):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = line
                        p.space_after = Pt(6)
                        for run in p.runs:
                            run.font.size = Pt(16)
                    break

        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = notes
        else:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    # Slide 1 — Title
    _add_slide(
        title="ONCORIX (rivolumab) — Field Team Training",
        body_lines=[],
        notes=(
            "Welcome to the ONCORIX training module. This deck covers the approved "
            "indication, key efficacy and safety data, and recommended talking points "
            "for HCP conversations."
        ),
        is_title_slide=True,
        subtitle="Unresectable/Metastatic Melanoma",
    )

    # Slide 2 — Indication & MOA
    _add_slide(
        title="Indication & Mechanism of Action",
        body_lines=[
            "ONCORIX is a humanised monoclonal antibody targeting PD-1.",
            "",
            "Approved indication:",
            "Treatment of adult patients with unresectable or metastatic melanoma.",
            "",
            "Mechanism: Blocks PD-1 receptor, re-activating the patient’s",
            "immune response against tumour cells.",
        ],
        notes=(
            "Key point: ONCORIX works by blocking the PD-1 receptor, re-activating "
            "the patient’s immune response against tumour cells. Only discuss the "
            "approved melanoma indication."
        ),
    )

    # Slide 3 — Efficacy
    _add_slide(
        title="MERIDIAN-301: Key Efficacy Results",
        body_lines=[
            "• Overall Survival: 24.1 vs 13.6 months (HR 0.68; p=0.002)",
            "• ORR: 42.3% vs 16.2%",
            "• Median PFS: 11.2 vs 5.4 months",
            "• 78% of responders maintained response at 12 months",
            "• Consistent benefit across PD-L1, ECOG, and prior therapy subgroups",
        ],
        notes=(
            "When presenting efficacy, always lead with the overall survival benefit. "
            "The ORR and PFS data support the OS finding. Remember: always pair efficacy "
            "with safety data (next slide)."
        ),
    )

    # Slide 4 — Safety
    _add_slide(
        title="Safety Profile",
        body_lines=[
            "• Most common AEs (≥20%): fatigue (38%), rash (27%), diarrhoea (24%), nausea (21%)",
            "• Immune-mediated AEs: 31% (mostly Grade 1-2, manageable)",
            "• Grade 3-4 treatment-related AEs: 18% vs 12% (control)",
            "• Discontinuation due to AEs: 9%",
            "• No new safety signals at 24-month follow-up",
        ],
        notes=(
            "Fair balance is mandatory. When an HCP asks about efficacy, always follow up "
            "with the safety profile. Emphasise that immune-mediated AEs are generally "
            "manageable with established protocols."
        ),
    )

    # Slide 5 — Dosing
    _add_slide(
        title="Dosing & Administration",
        body_lines=[
            "• Dose: 200 mg IV infusion over 30 minutes every 3 weeks",
            "• No dose adjustment for mild hepatic impairment",
            "• Continue until progression, unacceptable toxicity, or up to 24 months",
        ],
        notes=(
            "The Q3W dosing schedule is straightforward. Highlight the 30-minute infusion "
            "time as a practical advantage for patients and infusion centres."
        ),
    )

    # Slide 6 — Objection Handling
    _add_slide(
        title="Common HCP Objections & Approved Responses",
        body_lines=[
            'Objection: "What about the immune-related AE rate?"',
            'Response: "In MERIDIAN-301, immune-mediated AEs occurred in 31% of patients,',
            "but most were Grade 1-2 and manageable with established protocols.",
            'Only 9% of patients discontinued due to AEs."',
            "",
            'Objection: "How does this compare to existing checkpoint inhibitors?"',
            'Response: "I can share the MERIDIAN-301 data. For cross-trial comparisons,',
            "I'd recommend discussing with our Medical Science Liaison who can provide",
            'additional context."',
        ],
        notes=(
            "CRITICAL: Never make unsupported comparative claims. If asked to compare, "
            "share our trial data and refer to the MSL for further discussion. This is a "
            "compliance requirement."
        ),
    )

    prs.save(str(output_path))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Generate sample training fixture documents (PDF + PPTX).",
    )
    parser.add_argument(
        "--output-dir",
        default=str(DEFAULT_OUTPUT_DIR),
        help=f"Output directory (default: {DEFAULT_OUTPUT_DIR})",
    )
    args = parser.parse_args()

    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    errors: list[str] = []

    # PDF
    pdf_path = output_dir / "sample_product_message.pdf"
    try:
        _generate_pdf(pdf_path)
        print(f"  PDF: {pdf_path}")
    except ImportError:
        errors.append("reportlab is required for PDF generation. Install with: pip install reportlab")
    except Exception as e:
        errors.append(f"PDF generation failed: {e}")

    # PPTX
    pptx_path = output_dir / "sample_training_deck.pptx"
    try:
        _generate_pptx(pptx_path)
        print(f"  PPTX: {pptx_path}")
    except ImportError:
        errors.append("python-pptx is required for PPTX generation. Install with: pip install python-pptx")
    except Exception as e:
        errors.append(f"PPTX generation failed: {e}")

    if errors:
        for err in errors:
            print(f"  ERROR: {err}", file=sys.stderr)
        return 1

    print("Done.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
