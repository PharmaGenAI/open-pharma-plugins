"""ingest_document — extract structured text from a PDF or PPTX file."""

from __future__ import annotations

from typing import Any

from pydantic import BaseModel, Field


class IngestDocumentArgs(BaseModel):
    file_path: str = Field(description="Absolute path to the PDF or PPTX file to ingest")
    document_title: str | None = Field(
        default=None,
        description="Optional title override. If not provided, the title is extracted from the first page/slide.",
    )


TOOL: dict[str, Any] = {
    "name": "ingest_document",
    "description": (
        "Extract structured text from a PDF or PPTX file and store it in the "
        "training content index. Returns a summary with the document ID, title, "
        "page count, and file type. For PPTX files, extracts slide titles, body "
        "text, and speaker notes per slide. For PDF files, extracts the embedded "
        "text layer per page. Use this before search_content or get_document_page."
    ),
    "args": IngestDocumentArgs,
}


def handle(arguments: dict[str, Any]) -> list[dict[str, Any]]:
    import json
    import os
    from datetime import datetime, timezone

    from .._content_store import document_id_from_path, save_document

    file_path = arguments["file_path"]
    document_title = arguments.get("document_title")

    if not os.path.isfile(file_path):
        return [{"type": "text", "text": json.dumps({"error": f"File not found: {file_path}"})}]

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        pages, auto_title = _extract_pdf(file_path)
        file_type = "pdf"
    elif ext in (".pptx", ".ppt"):
        pages, auto_title = _extract_pptx(file_path)
        file_type = "pptx"
    else:
        return [{"type": "text", "text": json.dumps({"error": f"Unsupported file type: {ext}. Use .pdf or .pptx"})}]

    pages = [p for p in pages if p.get("text", "").strip()]

    title = document_title or auto_title or os.path.basename(file_path)
    doc_id = document_id_from_path(file_path)

    doc = {
        "document_id": doc_id,
        "file_name": os.path.basename(file_path),
        "file_path": os.path.abspath(file_path),
        "file_type": file_type,
        "title": title,
        "total_pages": len(pages),
        "pages": pages,
        "ingested_at": datetime.now(timezone.utc).isoformat(),
    }

    save_document(doc)

    summary = {
        "document_id": doc_id,
        "file_name": doc["file_name"],
        "file_type": file_type,
        "title": title,
        "total_pages": len(pages),
        "ingested_at": doc["ingested_at"],
    }
    return [{"type": "text", "text": json.dumps(summary, indent=2)}]


def _extract_pdf(file_path: str) -> tuple[list[dict], str | None]:
    import pypdfium2 as pdfium

    doc = pdfium.PdfDocument(file_path)
    pages: list[dict] = []
    auto_title: str | None = None

    for i in range(len(doc)):
        textpage = doc[i].get_textpage()
        text = textpage.get_text_bounded()
        textpage.close()
        pages.append(
            {
                "page_number": i + 1,
                "text": text,
                "page_type": "page",
                "slide_title": None,
                "speaker_notes": None,
            }
        )
        if i == 0 and text.strip():
            first_line = text.strip().split("\n")[0].strip()
            if first_line and len(first_line) < 200:
                auto_title = first_line

    doc.close()
    return pages, auto_title


def _extract_pptx(file_path: str) -> tuple[list[dict], str | None]:
    from pptx import Presentation

    prs = Presentation(file_path)
    pages: list[dict] = []
    auto_title: str | None = None

    for i, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text

        body_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                body_texts.append(shape.text_frame.text)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text

        pages.append(
            {
                "page_number": i + 1,
                "text": "\n".join(body_texts),
                "page_type": "slide",
                "slide_title": title or None,
                "speaker_notes": notes or None,
            }
        )

        if i == 0 and title:
            auto_title = title

    return pages, auto_title
